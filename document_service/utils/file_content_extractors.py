import PyPDF2
from pdf2image import convert_from_bytes
import os
import io
import tempfile
from docx import Document



import requests
import base64
from pptx import Presentation
import threading

API_BASE_URL = "https://llm.rdfai.ru/api/providers/openai/v1"
API_KEY = "9a49054fe339d4bf09f9a35f0d7169e4f4bfacd4bc4b9594406b3348c89d8558"


def encode_image(image_bytes):
    """Кодирует байты изображения в Base64."""
    return base64.b64encode(image_bytes).decode("utf-8")


def describe_image(image_bytes, image_name):
    """Отправляет запрос к API для описания изображения."""
    try:
        image_base64 = encode_image(image_bytes)
        payload = {
            "model": "gpt-4.1",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe the content of this image in a concise manner."},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{image_base64}"},
                        },
                    ],
                }
            ],
            "max_tokens": 100,
        }

        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }

        response = requests.post(
            f"{API_BASE_URL}/chat/completions",
            headers=headers,
            json=payload,
            timeout=30
        )
        response.raise_for_status()

        return response.json()['choices'][0]['message']['content']
    except Exception as e:
        return f"Error describing image {image_name}: {str(e)}"


def process_image(image, page_num, output_dir, extracted_content, lock):
    """Обрабатывает одно изображение в отдельном потоке."""
    image_ext = 'png'
    image_name = f"page_{page_num + 1}_img_1.{image_ext}"
    image_path = os.path.join(output_dir, image_name)

    # Сохранение изображения во временный файл
    image.save(image_path, 'PNG')

    # Чтение байтов изображения
    with open(image_path, "rb") as image_file:
        image_bytes = image_file.read()

    # Описание изображения через API
    description = describe_image(image_bytes, image_name)

    # Потокобезопасное добавление результата
    with lock:
        extracted_content.append({
            'type': 'image',
            'name': image_name,
            'description': description,
            'page': page_num + 1
        })


def extract_pdf_content(pdf_bytes, output_dir="extracted_images"):
    """Извлекает текст и описания изображений из PDF, переданного в виде байтов."""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Чтение PDF из байтов
    pdf_file = io.BytesIO(pdf_bytes)
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    extracted_content = []
    lock = threading.Lock()

    # Конвертация страниц PDF в изображения из байтов
    images = convert_from_bytes(pdf_bytes, thread_count=os.cpu_count())

    # Итерация по страницам
    threads = []
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]

        # Извлечение текста
        text = page.extract_text()
        if text.strip():
            with lock:
                extracted_content.append({
                    'type': 'text',
                    'content': text,
                    'page': page_num + 1
                })

        # Обработка изображений
        if page_num < len(images):
            thread = threading.Thread(
                target=process_image,
                args=(images[page_num], page_num, output_dir, extracted_content, lock)
            )
            threads.append(thread)
            thread.start()

    # Ожидание завершения всех потоков
    for thread in threads:
        thread.join()

    # Сортировка по номеру страницы
    extracted_content.sort(key=lambda x: x['page'])

    res = ""
    for item in extracted_content:
        if item['type'] == 'text':
            res += f"\n=== Page {item['page']} Text ===\n"
            res += item['content']
            res += "\n"
        elif item['type'] == 'image':
            res += f"\n=== Page {item['page']} Image ===\n"
            res += f"Image Name: {item['name']}\n"
            res += f"Description: {item['description']}\n"
            res += "\n"

    return res


def process_image_pptx(image_bytes, slide_num, image_num, output_dir, extracted_content, lock):
    """Обрабатывает одно изображение в отдельном потоке."""
    image_ext = 'png'
    image_name = f"slide_{slide_num + 1}_img_{image_num + 1}.{image_ext}"
    image_path = os.path.join(output_dir, image_name)

    # Сохранение изображения
    with open(image_path, "wb") as f:
        f.write(image_bytes)

    # Описание изображения через API
    description = describe_image(image_bytes, image_name)

    # Потокобезопасное добавление результата
    with lock:
        extracted_content.append({
            'type': 'image',
            'name': image_name,
            'description': description,
            'slide': slide_num + 1
        })

def extract_pptx_content_with_images(pptx_bytes, output_dir="extracted_pptx_images"):
    """Извлекает текст и описания изображений из PPTX, переданного в виде байтов."""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    extracted_content = []
    lock = threading.Lock()
    threads = []

    # Создание временного файла для PPTX
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        temp_file.write(pptx_bytes)
        temp_file_path = temp_file.name

    try:
        # Открытие PPTX
        prs = Presentation(temp_file_path)

        # Итерация по слайдам
        for slide_num, slide in enumerate(prs.slides):
            # Извлечение текста
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            if text.strip():
                with lock:
                    extracted_content.append({
                        'type': 'text',
                        'content': text.strip(),
                        'slide': slide_num + 1
                    })

            # Извлечение изображений
            image_num = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_bytes = shape.image.blob
                    image_num += 1

                    # Запуск обработки изображения в отдельном потоке
                    thread = threading.Thread(
                        target=process_image_pptx,
                        args=(image_bytes, slide_num, image_num, output_dir, extracted_content, lock)
                    )
                    threads.append(thread)
                    thread.start()

        # Ожидание завершения всех потоков
        for thread in threads:
            thread.join()

    finally:
        # Удаление временного файла
        os.unlink(temp_file_path)

    # Сортировка по номеру слайда
    extracted_content.sort(key=lambda x: x['slide'])

    res = ""
    for item in extracted_content:
        if item['type'] == 'text':
            res += f"\n=== Slide {item['slide']} Text ===\n"
            res += item['content']
            res += "\n"
        elif item['type'] == 'image':
            res += f"\n=== Slide {item['slide']} Image ===\n"
            res += f"Image Name: {item['name']}\n"
            res += f"Description: {item['description']}\n"
            res += "\n"
    return res


def process_image_docx(image_bytes, image_num, output_dir, extracted_content, lock):
    """Обрабатывает одно изображение в отдельном потоке."""
    image_ext = 'png'
    image_name = f"doc_img_{image_num + 1}.{image_ext}"
    image_path = os.path.join(output_dir, image_name)

    # Сохранение изображения
    with open(image_path, "wb") as f:
        f.write(image_bytes)

    # Описание изображения через API
    description = describe_image(image_bytes, image_name)

    # Потокобезопасное добавление результата
    with lock:
        extracted_content.append({
            'type': 'image',
            'name': image_name,
            'description': description
        })


def extract_docx_content_with_images(docx_bytes, output_dir="extracted_docx_images"):
    """Извлекает текст и описания изображений из DOCX, переданного в виде байтов."""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    extracted_content = []
    lock = threading.Lock()
    threads = []

    # Создание временного файла для DOCX
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
        temp_file.write(docx_bytes)
        temp_file_path = temp_file.name

    try:
        # Открытие DOCX
        doc = Document(temp_file_path)

        # Извлечение текста
        text = ""
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"

        # Извлечение текста из таблиц
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"

        if text.strip():
            extracted_content.append({
                'type': 'text',
                'content': text.strip()
            })

        # Извлечение изображений
        image_num = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_bytes = rel.target_part.blob
                image_num += 1

                # Запуск обработки изображения в отдельном потоке
                thread = threading.Thread(
                    target=process_image_docx,
                    args=(image_bytes, image_num, output_dir, extracted_content, lock)
                )
                threads.append(thread)
                thread.start()

        # Ожидание завершения всех потоков
        for thread in threads:
            thread.join()

    finally:
        # Удаление временного файла
        os.unlink(temp_file_path)

    res = ""
    for item in extracted_content:
        if item['type'] == 'text':
            res += "\n=== Document Text ===\n"
            res += item['content']
            res += "\n"
        elif item['type'] == 'image':
            res += f"\n=== Image {item['name']} ===\n"
            res += f"Image Name: {item['name']}\n"
            res += f"Description: {item['description']}\n"
            res += "\n"

    return res


def extract_pdf_text(file_content):
    full_text = ""
    with io.BytesIO(file_content) as pdf_file:  # noqa
        read_pdf = PyPDF2.PdfReader(pdf_file)
        for page in read_pdf.pages:
            full_text += page.extract_text()
    return full_text


def extract_pptx_text(pptx_bytes):
    """Извлекает текст из PPTX, переданного в виде байтов, и возвращает строку."""
    extracted_text = []

    # Создание временного файла для PPTX
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
        temp_file.write(pptx_bytes)
        temp_file_path = temp_file.name

    try:
        # Открытие PPTX
        prs = Presentation(temp_file_path)

        # Итерация по слайдам
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""

            # Извлечение текста из фигур
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"

            # Извлечение текста из заметок (если есть)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text += "Notes:\n" + notes + "\n"

            # Извлечение текста из таблиц
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text += cell.text + "\n"

            if slide_text.strip():
                extracted_text.append(f"=== Slide {slide_num + 1} Text ===\n{slide_text.strip()}")

    finally:
        # Удаление временного файла
        os.unlink(temp_file_path)

    # Объединение текста в одну строку с разделителями
    return "\n\n".join(extracted_text) if extracted_text else ""


def extract_docx_text(docx_bytes):
    """Извлекает текст из DOCX, переданного в виде байтов, и возвращает строку."""
    extracted_text = []

    # Создание временного файла для DOCX
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
        temp_file.write(docx_bytes)
        temp_file_path = temp_file.name

    try:
        # Открытие DOCX
        doc = Document(temp_file_path)

        # Извлечение текста из параграфов
        for para in doc.paragraphs:
            if para.text.strip():
                extracted_text.append(para.text.strip())

        # Извлечение текста из таблиц
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        extracted_text.append(cell.text.strip())

    finally:
        # Удаление временного файла
        os.unlink(temp_file_path)

    # Объединение текста в одну строку с разделителями
    return "\n".join(extracted_text) if extracted_text else ""
